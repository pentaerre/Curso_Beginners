from pptx import Presentation
from pptx.util import Inches
img_path = r"C:\Users\U40687\Desktop\PPT_EXAMPLE\INPUT\PYTHON_ANIMAL.PNG"

# Lo primero que se hace es un objeto presentacion
root = Presentation()

# Elegimos un tipo de slide, debajo se explica los tipos que hay
first_slide_layout = root.slide_layouts[0]

""" Ref for slide types:  
0 ->  title and subtitle 
1 ->  title and content 
2 ->  section header 
3 ->  two content 
4 ->  Comparison 
5 ->  Title only  
6 ->  Blank 
7 ->  Content with caption 
8 ->  Pic with caption 
"""

# Se añade una slide del tipo que hemos definido que es de titulo
slide = root.slides.add_slide(first_slide_layout)

# Adding title and subtitle in  
# slide i.e. first page of slide  
slide.shapes.title.text = "Creador de diapositivas Altran"

# We have different formats of  
# subtitles in ppts, for simple 
# subtitle this method should  
# implemented, you can change 
# 0 to 1 for different design 
slide.placeholders[1].text = " This is 2nd way"


# Selecting blank slide
blank_slide_layout = root.slide_layouts[6]
# Attaching slide to ppt
slide = root.slides.add_slide(blank_slide_layout)

# For margins
left = top = Inches(1)

# adding images
pic = slide.shapes.add_picture(img_path,
                               left, top)

left = Inches(1)
height = Inches(1)

pic = slide.shapes.add_picture(img_path, left,
                               top, height=height)


# Saving file 
root.save(r"C:\Users\U40687\Desktop\PPT_EXAMPLE\OUTPUT\Output.pptx")

print("done")